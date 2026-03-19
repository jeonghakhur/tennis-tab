"""Tennis-Tab 소개 PPT — 사용자 친화적 버전 (흰색 배경, 한글 중심)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 색상
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1F, 0x20, 0x37)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
SUB_TEXT = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x10, 0xB9, 0x81)
BLUE = RGBColor(0x38, 0x82, 0xF6)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
YELLOW_BG = RGBColor(0xFF, 0xF7, 0xED)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
BORDER_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = "맑은 고딕"


def bg(slide, color=WHITE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, radius=None, border_color=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.adjustments[0] = radius
    return s


def txt(slide, l, t, w, h, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb


def multiline(slide, l, t, w, h, lines, size=13, color=DARK_TEXT, line_spacing=1.3):
    """여러 줄 텍스트를 하나의 텍스트박스에"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(size * (line_spacing - 1) + 2)
    return tb


def card(slide, l, t, w, h, icon, title, items, accent=GREEN, bg_color=None):
    """아이콘 + 제목 + 항목 리스트 카드"""
    card_bg = bg_color or CARD_BG
    rect(slide, l, t, w, h, card_bg, radius=0.04, border_color=BORDER_LIGHT)
    # 상단 악센트 라인
    rect(slide, l + Inches(0.3), t + Inches(0.15), Inches(0.5), Inches(0.04), accent, radius=0.5)

    txt(slide, l + Inches(0.3), t + Inches(0.35), w - Inches(0.6), Inches(0.45),
        f"{icon}  {title}", size=16, color=DARK_TEXT, bold=True)

    item_lines = [f"•  {item}" for item in items]
    multiline(slide, l + Inches(0.35), t + Inches(0.85), w - Inches(0.7), h - Inches(1.0),
              item_lines, size=12, color=SUB_TEXT)


def pill(slide, l, t, text, accent=GREEN, bg_color=None):
    """둥근 배지/필"""
    pill_bg = bg_color or accent
    w = Inches(len(text) * 0.115 + 0.5)
    s = rect(slide, l, t, w, Inches(0.38), pill_bg, radius=0.4)
    tf = s.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE if bg_color is None else accent
    p.font.bold = True
    p.font.name = FONT
    return w


def flow_step(slide, x, y, text, accent=GREEN, bg_c=None):
    """흐름도 단계 (둥근 사각형)"""
    w = Inches(len(text) * 0.115 + 0.6)
    fill = bg_c or accent
    s = rect(slide, x, y, w, Inches(0.5), fill, radius=0.15)
    tf = s.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    return w


def arrow(slide, x, y):
    txt(slide, x, y - Inches(0.02), Inches(0.35), Inches(0.5), "→", size=16, color=SUB_TEXT,
        align=PP_ALIGN.CENTER)


# =====================================================
# 슬라이드 1: 표지
# =====================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1, WHITE)

# 상단 그린 바
rect(s1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

# 중앙 콘텐츠
txt(s1, Inches(1), Inches(2.2), Inches(11), Inches(0.5),
    "테니스를 더 쉽고, 더 즐겁게", size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

txt(s1, Inches(1), Inches(2.9), Inches(11), Inches(1.2),
    "Tennis-Tab", size=56, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

txt(s1, Inches(2), Inches(4.3), Inches(9), Inches(0.8),
    "대회 참가부터 클럽 활동, 레슨 관리까지\n테니스에 필요한 모든 것을 한곳에서",
    size=18, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# 핵심 키워드 필
keywords = ["대회 관리", "클럽 운영", "레슨", "커뮤니티", "AI 도우미"]
kw_colors = [GREEN, BLUE, ORANGE, PURPLE, RGBColor(0xF5, 0x9E, 0x0B)]
x = Inches(3.0)
for kw, kc in zip(keywords, kw_colors):
    pill(s1, x, Inches(5.5), kw, kc)
    x += Inches(len(kw) * 0.115 + 0.65)


# =====================================================
# 슬라이드 2: 한눈에 보는 주요 기능
# =====================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s2, WHITE)

txt(s2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "한눈에 보는 주요 기능", size=30, color=BLACK, bold=True)
rect(s2, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), GREEN)

features = [
    ("🏆", "대회 관리", [
        "대회 개설부터 마감까지 전 과정 관리",
        "단식 · 복식 · 단체전 모두 지원",
        "자동 대진표 생성 (조편성 → 토너먼트)",
        "참가 신청 & 온라인 결제",
    ], GREEN, GREEN_BG),
    ("🎾", "클럽 운영", [
        "우리 클럽만의 공간 생성",
        "회원 가입 · 관리 · 역할 설정",
        "정기 모임 일정 & 참석 관리",
        "모임별 자동 대진 & 경기 결과 기록",
    ], BLUE, BLUE_BG),
    ("📚", "레슨 시스템", [
        "코치 프로필 & 레슨 프로그램 등록",
        "수준별 맞춤 레슨 (입문~고급)",
        "수강 신청 & 일정 예약",
        "출석 관리 & 월별 결제 기록",
    ], ORANGE, ORANGE_BG),
]

x = Inches(0.8)
for icon, title, items, accent, bg_c in features:
    card(s2, x, Inches(1.5), Inches(3.7), Inches(2.7), icon, title, items, accent, bg_c)
    x += Inches(3.9)

bottom = [
    ("💬", "커뮤니티", [
        "공지 · 자유 · 정보 · 대회후기 게시판",
        "사진 · 파일 첨부, 댓글, 좋아요",
    ], PURPLE, PURPLE_BG),
    ("🤖", "AI 도우미", [
        "대화로 대회 검색 & 참가 신청",
        "궁금한 점은 챗봇에게 바로 질문",
    ], RGBColor(0xF5, 0x9E, 0x0B), YELLOW_BG),
    ("🔔", "알림", [
        "참가 승인 · 대진표 발표 등 실시간 알림",
        "카카오 알림톡으로도 받아보기",
    ], RED, RED_BG),
]

x = Inches(0.8)
for icon, title, items, accent, bg_c in bottom:
    card(s2, x, Inches(4.5), Inches(3.7), Inches(2.2), icon, title, items, accent, bg_c)
    x += Inches(3.9)


# =====================================================
# 슬라이드 3: 대회 — 개설부터 시상까지
# =====================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s3, WHITE)

txt(s3, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "🏆  대회 — 개설부터 시상까지", size=30, color=BLACK, bold=True)
rect(s3, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), GREEN)

# 대회 흐름
txt(s3, Inches(0.8), Inches(1.4), Inches(10), Inches(0.4),
    "대회 진행 흐름", size=16, color=GREEN, bold=True)

steps = [
    ("대회 개설", GREEN),
    ("참가 모집", BLUE),
    ("모집 마감", ORANGE),
    ("대진표 생성", PURPLE),
    ("경기 진행", RGBColor(0xF5, 0x9E, 0x0B)),
    ("시상 & 완료", RED),
]
x = Inches(0.8)
for i, (step, clr) in enumerate(steps):
    w = flow_step(s3, x, Inches(1.9), step, clr)
    if i < len(steps) - 1:
        arrow(s3, x + w + Inches(0.02), Inches(1.9))
    x += w + Inches(0.4)

# 참가 신청 흐름
txt(s3, Inches(0.8), Inches(2.7), Inches(10), Inches(0.4),
    "참가 신청 과정", size=16, color=GREEN, bold=True)

entry_steps = [
    ("대회 선택", BLUE),
    ("부서 · 종목 선택", PURPLE),
    ("선수 정보 입력", ORANGE),
    ("참가비 결제", GREEN),
    ("신청 완료!", RGBColor(0xF5, 0x9E, 0x0B)),
]
x = Inches(0.8)
for i, (step, clr) in enumerate(entry_steps):
    w = flow_step(s3, x, Inches(3.2), step, clr)
    if i < len(entry_steps) - 1:
        arrow(s3, x + w + Inches(0.02), Inches(3.2))
    x += w + Inches(0.4)

# 하단 카드
card(s3, Inches(0.8), Inches(4.1), Inches(5.8), Inches(3.0),
     "📊", "대진표 자동 생성", [
         "드래그 & 드롭으로 조편성",
         "예선 조별리그 → 본선 토너먼트 자동 연결",
         "128강부터 결승까지 자동 배치",
         "3 · 4위전 선택 가능",
         "경기 결과 입력 시 다음 라운드 자동 반영",
     ], GREEN, GREEN_BG)

card(s3, Inches(7.0), Inches(4.1), Inches(5.8), Inches(3.0),
     "🎾", "다양한 경기 형식", [
         "개인전 단식 — 1:1 경기",
         "개인전 복식 — 파트너와 함께 참가",
         "단체전 단식 — 팀 대표 선수들의 개인 경기",
         "단체전 복식 — 팀 대표 조합의 복식 경기",
         "세트별 출전 선수 배정 & 점수 기록",
     ], BLUE, BLUE_BG)


# =====================================================
# 슬라이드 4: 클럽 운영
# =====================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s4, WHITE)

txt(s4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "🎾  클럽 — 우리만의 테니스 공간", size=30, color=BLACK, bold=True)
rect(s4, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), BLUE)

card(s4, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.7),
     "🏠", "클럽 만들기 & 회원 관리", [
         "누구나 클럽을 만들고 회원을 초대할 수 있어요",
         "자유 가입 · 승인제 · 초대 전용 중 선택",
         "운영자 · 관리자 · 일반 회원 역할 구분",
         "지역별로 가까운 클럽 검색",
         "소속 협회와 연결 가능",
     ], BLUE, BLUE_BG)

card(s4, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.7),
     "📅", "정기 모임 관리", [
         "모임 일정 등록 & 참석 여부 확인",
         "성별 · 시간대 고려한 자동 대진 편성",
         "단식 · 남복 · 여복 · 혼복 모두 지원",
         "외부 게스트도 참여 가능",
         "점수 이견 시 관리자가 최종 판정",
     ], GREEN, GREEN_BG)

card(s4, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.7),
     "📊", "경기 기록 & 통계", [
         "모임별 경기 결과 자동 저장",
         "회원별 승률 · 전적 · 출석률 확인",
         "최근 경기 이력 한눈에 보기",
         "클럽 내 랭킹으로 재미있는 경쟁",
     ], PURPLE, PURPLE_BG)

card(s4, Inches(7.0), Inches(4.4), Inches(5.8), Inches(2.7),
     "👤", "회원 프로필", [
         "개인 전적 페이지에서 통산 기록 확인",
         "소속 클럽별 경기 이력 분리 표시",
         "수상 기록 모아보기",
         "내 활동 내역 한곳에서 관리",
     ], ORANGE, ORANGE_BG)


# =====================================================
# 슬라이드 5: 레슨
# =====================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s5, WHITE)

txt(s5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "📚  레슨 — 체계적인 테니스 교육", size=30, color=BLACK, bold=True)
rect(s5, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), ORANGE)

card(s5, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.7),
     "👨‍🏫", "코치 & 프로그램", [
         "코치 프로필 등록 (경력, 자격증, 사진)",
         "수준별 프로그램: 입문 · 초급 · 중급 · 고급",
         "평일 / 주말, 1인 / 2인 요금 차등 설정",
         "프로그램 공개 · 마감 상태 관리",
         "클럽 소속 코치 연결",
     ], ORANGE, ORANGE_BG)

card(s5, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.7),
     "📝", "수강 신청 & 예약", [
         "원하는 시간대에 레슨 신청",
         "빈 자리가 없으면 대기 등록 가능",
         "레슨 일정은 캘린더로 한눈에 확인",
         "일정 변경 요청 & 승인 프로세스",
         "빈 슬롯이 없을 때 희망 일정 문의",
     ], BLUE, BLUE_BG)

card(s5, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.7),
     "✅", "출석 & 결제 관리", [
         "레슨별 출석 · 결석 · 지각 기록",
         "월별 결제 내역 관리 (계좌이체 · 현금)",
         "수강 연장 시 카카오 알림톡 안내",
         "코치 전용 관리 페이지 제공",
     ], GREEN, GREEN_BG)

card(s5, Inches(7.0), Inches(4.4), Inches(5.8), Inches(2.7),
     "🔍", "레슨 검색", [
         "지역 · 코치 · 수준 · 요일 · 시간으로 필터",
         "클럽 레슨 프로그램 직접 조회",
         "내가 수강 중인 레슨 모아보기",
         "레슨 관련 문의 & 답변",
     ], PURPLE, PURPLE_BG)


# =====================================================
# 슬라이드 6: 편의 기능
# =====================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s6, WHITE)

txt(s6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "더 편리하게 사용하는 방법", size=30, color=BLACK, bold=True)
rect(s6, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), GREEN)

card(s6, Inches(0.8), Inches(1.4), Inches(3.7), Inches(2.8),
     "🤖", "AI 도우미", [
         '"이번 달 대회 뭐 있어?" 물어보면 검색',
         '"복식 참가 신청할래" 하면 안내 시작',
         '파트너 · 팀원 등록도 대화로 해결',
         '참가 취소도 간편하게',
         '대진표 · 수상 기록 조회',
     ], RGBColor(0xF5, 0x9E, 0x0B), YELLOW_BG)

card(s6, Inches(4.8), Inches(1.4), Inches(3.7), Inches(2.8),
     "🔔", "알림 센터", [
         "참가 승인 · 거절 알림",
         "대진표 발표 · 경기 결과 알림",
         "클럽 가입 승인 · 초대 알림",
         "환불 완료 · 문의 답변 알림",
         "카카오 알림톡으로도 수신 가능",
     ], RED, RED_BG)

card(s6, Inches(8.8), Inches(1.4), Inches(3.7), Inches(2.8),
     "💳", "간편 결제", [
         "대회 참가비 온라인 결제",
         "결제 상태 실시간 확인",
         "관리자 환불 처리 지원",
         "레슨비 결제 기록 관리",
     ], GREEN, GREEN_BG)

card(s6, Inches(0.8), Inches(4.5), Inches(3.7), Inches(2.5),
     "👤", "마이페이지", [
         "내 프로필 & 대회 참가 내역",
         "소속 클럽 & 회원 전적",
         "수강 중인 레슨 현황",
         "알림 모아보기",
     ], BLUE, BLUE_BG)

card(s6, Inches(4.8), Inches(4.5), Inches(3.7), Inches(2.5),
     "💬", "커뮤니티", [
         "공지 · 자유 · 정보 · 대회후기",
         "사진 첨부 & 리치 텍스트 작성",
         "댓글 · 좋아요로 소통",
         "카카오톡으로 게시글 공유",
     ], PURPLE, PURPLE_BG)

card(s6, Inches(8.8), Inches(4.5), Inches(3.7), Inches(2.5),
     "🏅", "수상 기록", [
         "대회별 수상 내역 자동 저장",
         "우승 · 준우승 · 3위 기록",
         "선수별 · 클럽별 수상 조회",
         "개인 프로필에 자동 연결",
     ], ORANGE, ORANGE_BG)


# =====================================================
# 슬라이드 7: 관리자 기능
# =====================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s7, WHITE)

txt(s7, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
    "⚙️  관리자를 위한 도구", size=30, color=BLACK, bold=True)
rect(s7, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.05), GREEN)

# 권한 설명
txt(s7, Inches(0.8), Inches(1.4), Inches(10), Inches(0.4),
    "누가 무엇을 할 수 있나요?", size=16, color=GREEN, bold=True)

roles = [
    ("최고 관리자", "시스템 전체 관리, 다른 관리자 지정", RED, RED_BG),
    ("관리자", "대회 개설, 참가자 심사, 대진표 관리", ORANGE, ORANGE_BG),
    ("운영자", "자신이 만든 대회만 관리", BLUE, BLUE_BG),
    ("일반 회원", "대회 참가, 클럽 활동, 레슨 수강", GREEN, GREEN_BG),
]

x = Inches(0.8)
for name, desc, accent, bg_c in roles:
    rect(s7, x, Inches(1.9), Inches(2.9), Inches(1.3), bg_c, radius=0.04, border_color=BORDER_LIGHT)
    rect(s7, x + Inches(0.2), Inches(2.05), Inches(0.4), Inches(0.04), accent, radius=0.5)
    txt(s7, x + Inches(0.2), Inches(2.2), Inches(2.5), Inches(0.35),
        name, size=15, color=DARK_TEXT, bold=True)
    txt(s7, x + Inches(0.2), Inches(2.6), Inches(2.5), Inches(0.5),
        desc, size=11, color=SUB_TEXT)
    x += Inches(3.05)

# 관리 도구
admin_cards = [
    ("🏆", "대회 관리", [
        "대회 생성 · 수정 · 상태 변경",
        "부서별 설정 & 참가자 심사",
        "대진표 생성 & 결과 입력",
        "수상 기록 관리",
    ], GREEN, GREEN_BG),
    ("🎾", "클럽 · 레슨 관리", [
        "클럽 생성 & 회원 관리",
        "모임 일정 & 경기 결과",
        "코치 등록 & 레슨 프로그램",
        "수강생 출석 & 결제 관리",
    ], BLUE, BLUE_BG),
    ("📋", "운영 지원", [
        "협회 등록 & 담당자 관리",
        "고객 문의 접수 & 답변",
        "FAQ & 사용 가이드 관리",
        "공지사항 & 알림 발송",
    ], PURPLE, PURPLE_BG),
    ("💰", "결제 · 환불", [
        "참가비 결제 상태 확인",
        "환불 요청 처리",
        "레슨비 결제 현황",
        "결제 내역 조회",
    ], ORANGE, ORANGE_BG),
]

x = Inches(0.8)
for icon, title, items, accent, bg_c in admin_cards:
    card(s7, x, Inches(3.6), Inches(2.9), Inches(2.9), icon, title, items, accent, bg_c)
    x += Inches(3.05)


# =====================================================
# 슬라이드 8: 마무리
# =====================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s8, WHITE)

rect(s8, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GREEN)

txt(s8, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
    "테니스를 더 쉽고, 더 즐겁게", size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

txt(s8, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
    "Tennis-Tab", size=52, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

txt(s8, Inches(2), Inches(3.3), Inches(9), Inches(0.8),
    "대회 관리부터 클럽 운영, 레슨, 커뮤니티까지\n테니스 생활의 모든 순간을 함께합니다",
    size=18, color=SUB_TEXT, align=PP_ALIGN.CENTER)

# 핵심 수치
stats = [
    ("6가지+", "핵심 기능", "대회 · 클럽 · 레슨\n커뮤니티 · 알림 · 결제"),
    ("다양한", "경기 형식", "단식 · 복식 · 단체전\n예선 · 본선 자동 연결"),
    ("실시간", "알림 시스템", "카카오 알림톡\n인앱 알림 센터"),
    ("AI", "스마트 도우미", "대화로 대회 검색\n참가 신청까지 한번에"),
]

x = Inches(1.2)
for num, title, desc in stats:
    rect(s8, x, Inches(4.3), Inches(2.6), Inches(2.2), CARD_BG, radius=0.04, border_color=BORDER_LIGHT)
    rect(s8, x + Inches(0.9), Inches(4.4), Inches(0.8), Inches(0.04), GREEN, radius=0.5)
    txt(s8, x, Inches(4.55), Inches(2.6), Inches(0.5),
        num, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    txt(s8, x, Inches(5.1), Inches(2.6), Inches(0.4),
        title, size=15, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    txt(s8, x, Inches(5.5), Inches(2.6), Inches(0.8),
        desc, size=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)
    x += Inches(2.8)


# 저장
output = "/home/user/tennis-tab/Tennis-Tab_소개.pptx"
prs.save(output)
print(f"PPT 생성 완료: {output}")
